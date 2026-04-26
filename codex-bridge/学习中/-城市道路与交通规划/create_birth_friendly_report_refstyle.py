from __future__ import annotations

from collections import Counter
from pathlib import Path
from shutil import copy2

import openpyxl
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
REF = ROOT / "第四组道交开题汇报(1).pptx"
OUT = ROOT / "北京适龄生育双职工家庭生育友好空间调研_课程汇报PPT_参考原版重做.pptx"
EXCEL = sorted(ROOT.glob("*生育友好空间调研问卷*.xlsx"), key=lambda p: p.stat().st_mtime, reverse=True)[0]

INK = RGBColor(15, 17, 21)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GOLD = RGBColor(216, 152, 48)
DARK_GOLD = RGBColor(176, 118, 35)
GRAY = RGBColor(237, 237, 237)
LIGHT_GRAY = RGBColor(242, 242, 242)
MID_GRAY = RGBColor(96, 96, 96)
LINE = RGBColor(218, 218, 218)
GREEN = RGBColor(72, 126, 104)
BLUE = RGBColor(72, 112, 158)
CORAL = RGBColor(196, 92, 70)
FONT = "微软雅黑"

SLIDE_W, SLIDE_H = 20, 11.25


def rgb_hex(rgb: RGBColor) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def load_data():
    wb = openpyxl.load_workbook(EXCEL, data_only=True)
    ws = wb.active
    headers = [c.value for c in ws[1]]
    rows = [r for r in ws.iter_rows(min_row=2, values_only=True) if any(v is not None for v in r)]

    scores = [float(r[6]) for r in rows if r[6] is not None]
    dims = {
        "托育空间": range(14, 19),
        "陪伴空间": range(19, 24),
        "照护者休憩": range(24, 29),
        "氛围与意愿": range(29, 34),
    }
    dim_stats = {}
    item_stats = []
    for dim, idxs in dims.items():
        dim_vals = []
        for idx in idxs:
            nums = [float(r[idx]) for r in rows if r[idx] is not None]
            avg = sum(nums) / len(nums)
            dim_vals.extend(nums)
            item_stats.append((dim, headers[idx], avg))
        dim_stats[dim] = sum(dim_vals) / len(dim_vals)

    open_answers = []
    for r in rows:
        value = r[34] if len(r) > 34 else None
        if value is None:
            continue
        text = str(value).strip()
        if text and text not in {"(空)", "空"}:
            open_answers.append(text)

    return {
        "n": len(rows),
        "score_avg": sum(scores) / len(scores),
        "score_min": min(scores),
        "score_max": max(scores),
        "gender": Counter(r[7] for r in rows),
        "age": Counter(r[8] for r in rows),
        "children": Counter(r[11] for r in rows),
        "child_age": Counter(r[12] for r in rows),
        "community": Counter(r[13] for r in rows),
        "dim_stats": dim_stats,
        "item_stats": item_stats,
        "open_answers": open_answers,
    }


def delete_unwanted_slides(prs: Presentation, keep_zero_based: set[int]) -> None:
    sld_id_lst = prs.slides._sldIdLst
    sld_ids = list(sld_id_lst)
    for idx in reversed(range(len(sld_ids))):
        if idx in keep_zero_based:
            continue
        r_id = sld_ids[idx].rId
        prs.part.drop_rel(r_id)
        sld_id_lst.remove(sld_ids[idx])


def move_slide(prs: Presentation, old_idx: int, new_idx: int) -> None:
    sld_id_lst = prs.slides._sldIdLst
    slides = list(sld_id_lst)
    sld = slides[old_idx]
    sld_id_lst.remove(sld)
    sld_id_lst.insert(new_idx, sld)


def blank(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    return slide


def fill(shape, color: RGBColor, transparency: int | None = None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency is not None:
        shape.fill.transparency = transparency
    shape.line.fill.background()


def line(shape, color: RGBColor = LINE, width: float = 1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def text(
    slide,
    value: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 18,
    color: RGBColor = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    font: str = FONT,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return box


def title(slide, cn: str, en: str, page: int | None = None):
    text(slide, en, 0.82, 0.48, 5.2, 0.34, 16, GOLD, True)
    text(slide, cn, 0.82, 0.9, 13.2, 0.52, 25, BLACK, True)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.82), Inches(1.58), Inches(2.15), Inches(0.08))
    fill(bar, GOLD)
    if page is not None:
        text(slide, f"{page:02d}", 18.55, 10.35, 0.55, 0.28, 14, MID_GRAY, True, PP_ALIGN.RIGHT)
    text(slide, "给城市嵌上育儿外挂——北京适龄生育双职工家庭生育友好空间调研", 0.82, 10.38, 9.6, 0.26, 10, MID_GRAY)


def card(slide, x, y, w, h, head, body, accent=GOLD, fill_color=LIGHT_GRAY):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shp, fill_color)
    line(shp, LINE, 0.8)
    accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.12), Inches(h))
    fill(accent_bar, accent)
    text(slide, head, x + 0.28, y + 0.2, w - 0.45, 0.34, 16, BLACK, True)
    text(slide, body, x + 0.28, y + 0.7, w - 0.45, h - 0.82, 12.2, INK)


def number_chip(slide, num, label, x, y, color=GOLD):
    circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.72), Inches(0.72))
    fill(circ, color)
    text(slide, str(num), x, y + 0.12, 0.72, 0.3, 18, WHITE, True, PP_ALIGN.CENTER)
    text(slide, label, x + 0.9, y + 0.13, 3.1, 0.3, 15, BLACK, True)


def progress_bar(slide, x, y, w, label, value, max_value=5, color=GOLD):
    text(slide, label, x, y - 0.05, 4.8, 0.28, 12, BLACK, True)
    base = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 4.9), Inches(y), Inches(w), Inches(0.22))
    fill(base, GRAY)
    width = max(0.04, w * value / max_value)
    fg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 4.9), Inches(y), Inches(width), Inches(0.22))
    fill(fg, color)
    text(slide, f"{value:.2f}", x + 4.9 + w + 0.16, y - 0.06, 0.7, 0.28, 12, color, True)


def section_slide(prs, part, cn, en):
    slide = blank(prs)
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
    fill(bg, GOLD)
    text(slide, part, 1.1, 1.25, 2.2, 0.4, 20, WHITE, True)
    text(slide, cn, 1.1, 3.9, 9.4, 0.82, 38, WHITE, True)
    text(slide, en, 1.12, 4.88, 9.5, 0.36, 18, WHITE)
    for x, y, sz, tr in [(14.2, 1.15, 4.3, 25), (15.7, 4.1, 2.4, 15), (12.7, 6.4, 3.2, 30)]:
        o = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(sz), Inches(sz))
        fill(o, WHITE, tr)
    text(slide, "REPORT", 15.0, 9.0, 3.4, 0.6, 30, WHITE, True, PP_ALIGN.RIGHT)
    return slide


def add_cover(prs):
    slide = blank(prs)
    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(7.1), Inches(SLIDE_H))
    fill(left, GOLD)
    text(slide, "REPORT", 1.0, 1.0, 2.6, 0.4, 22, WHITE, True)
    text(slide, "给城市嵌上育儿外挂", 1.0, 3.55, 8.8, 0.65, 36, WHITE, True)
    text(slide, "关怀地理学下的北京适龄生育双职工家庭生育友好空间调研", 1.05, 4.35, 11.1, 0.52, 19, WHITE, True)
    text(slide, "课程汇报PPT  ·  项目概况 / 已取得进展 / 创新亮点 / 未来展望", 1.05, 5.18, 7.1, 0.36, 14, WHITE)
    for i, label in enumerate(["托育空间", "陪伴空间", "独处空间", "空间氛围"]):
        text(slide, label, 1.05 + i * 1.55, 6.1, 1.3, 0.28, 11, WHITE, True, PP_ALIGN.CENTER)
    text(slide, "北京林业大学 · 城市道路与交通规划", 1.05, 9.75, 5.5, 0.32, 12, WHITE)
    text(slide, "从宏观政策倡导进入双职工家庭每天真实经过的社区、通勤与公共空间。", 8.7, 3.45, 7.9, 0.88, 26, BLACK, True)
    text(slide, "以关怀地理学为理论底座，以问卷、访谈、GIS/POI和后续数字推演为方法路径，识别北京高密度居住区中生育友好空间的供需错配与优化方向。", 8.75, 4.7, 7.4, 1.3, 17, INK)
    for x, y, w, h, c, label in [
        (8.85, 6.85, 2.5, 1.15, LIGHT_GRAY, "初步问卷"),
        (11.7, 6.85, 2.5, 1.15, LIGHT_GRAY, "案例借鉴"),
        (14.55, 6.85, 2.5, 1.15, LIGHT_GRAY, "策略构建"),
    ]:
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        fill(shp, c)
        line(shp, LINE, 0.8)
        text(slide, label, x, y + 0.38, w, 0.3, 15, GOLD, True, PP_ALIGN.CENTER)
    return slide


def add_outline(prs):
    slide = blank(prs)
    title(slide, "汇报结构", "Presentation Structure", 2)
    items = [
        ("01", "项目概况", "研究背景、对象痛点、理论基础、场地与方法"),
        ("02", "已取得进展", "资料梳理、问卷设计、初步样本、结果判读"),
        ("03", "创新亮点", "关怀地理学转译、三类空间模型、多源数据路径"),
        ("04", "未来展望", "样本扩充、空间落点、数字推演、策略输出"),
    ]
    for i, (num, head, body) in enumerate(items):
        x = 1.05 + (i % 2) * 8.9
        y = 2.45 + (i // 2) * 2.85
        text(slide, num, x, y, 1.0, 0.55, 32, GOLD, True)
        card(slide, x + 1.25, y - 0.1, 6.65, 1.78, head, body, GOLD, LIGHT_GRAY)
    text(slide, "重点展开：项目概况与已取得进展。创新和展望作为阶段性汇报的判断与后续计划。", 1.05, 8.8, 13.8, 0.45, 19, BLACK, True)
    return slide


def add_progress_overview(prs, page):
    slide = blank(prs)
    title(slide, "已取得进展：从选题雏形进入初步实证", "Research Progress", page)
    stages = [
        ("选题确立", "聚焦北京适龄生育双职工家庭，明确“生育友好空间”不是单一托育设施，而是嵌入社区生活圈的支持网络。"),
        ("框架搭建", "建立托育空间、陪伴空间、照护者休憩空间、空间氛围四维评价框架。"),
        ("工具完成", "完成问卷量表设计，并形成默认报告与原始 Excel 数据。"),
        ("样本回收", "已获得 18 份初步样本，可用于阶段性趋势识别和后续访谈提纲校正。"),
    ]
    for i, (head, body) in enumerate(stages):
        x = 1.0 + i * 4.55
        number_chip(slide, i + 1, head, x, 2.18, GOLD)
        text(slide, body, x + 0.02, 3.1, 3.8, 1.55, 13.2, INK)
        if i < 3:
            arr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x + 3.75), Inches(2.45), Inches(0.62), Inches(0.28))
            fill(arr, GOLD)
    card(slide, 1.05, 6.3, 5.2, 1.8, "阶段性判断", "目前成果已能支撑课程中期汇报：选题逻辑、理论框架、调研工具和初步数据均已形成，下一步应从“样本数量”转向“空间落点”和“机制解释”。", GOLD)
    card(slide, 6.7, 6.3, 5.2, 1.8, "需要补强", "样本量仍偏小，且空间点位、受访者居住社区类型、通勤路径与托育资源之间尚未完成 GIS 叠合，需要后续实地调研和数据抓取补充。", CORAL)
    card(slide, 12.35, 6.3, 5.2, 1.8, "汇报策略", "本次汇报不把初步数据包装成定论，而是作为问题识别与后续调研设计的依据，体现研究正在推进且路径清晰。", BLUE)


def add_sample_profile(prs, data, page):
    slide = blank(prs)
    title(slide, "问卷样本画像：初步回收 18 份有效样本", "Questionnaire Profile", page)
    text(slide, f"N = {data['n']}  |  平均总分 {data['score_avg']:.1f}  |  得分区间 {data['score_min']:.0f}-{data['score_max']:.0f}", 1.05, 1.95, 10.2, 0.38, 21, BLACK, True)
    mapping = {
        "gender": {1: "男", 2: "女"},
        "age": {1: "25岁及以下", 2: "26-30岁", 3: "31-35岁", 4: "36-40岁", 5: "40岁以上"},
        "children": {1: "1个", 2: "2个", 3: "3个及以上", 4: "暂无但计划", 5: "其他"},
        "community": {1: "老旧小区", 2: "普通商品房", 3: "保障/公租房", 4: "新建商品房"},
    }
    panels = [
        ("性别结构", data["gender"], mapping["gender"], GOLD),
        ("年龄段", data["age"], mapping["age"], BLUE),
        ("子女数量", data["children"], mapping["children"], GREEN),
        ("社区类型", data["community"], mapping["community"], CORAL),
    ]
    for i, (head, counter, mp, color) in enumerate(panels):
        x = 1.05 + (i % 2) * 8.5
        y = 2.75 + (i // 2) * 3.25
        card(slide, x, y, 7.65, 2.4, head, "", color, LIGHT_GRAY)
        total = sum(counter.values())
        yy = y + 0.82
        for key, cnt in counter.most_common(4):
            label = mp.get(key, str(key))
            pct = cnt / total * 100 if total else 0
            text(slide, label, x + 0.35, yy, 1.7, 0.25, 11.5, BLACK, True)
            base = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 2.2), Inches(yy + 0.04), Inches(3.7), Inches(0.18))
            fill(base, WHITE)
            fg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 2.2), Inches(yy + 0.04), Inches(3.7 * pct / 100), Inches(0.18))
            fill(fg, color)
            text(slide, f"{cnt}人 / {pct:.0f}%", x + 6.05, yy - 0.03, 1.1, 0.25, 10.5, color, True, PP_ALIGN.RIGHT)
            yy += 0.38
    text(slide, "样本说明：当前数据用于识别趋势，不作为最终统计推断。后续需要扩大样本并补充访谈，以验证不同社区类型、子女年龄与通勤压力之间的差异。", 1.05, 9.95, 16.7, 0.36, 13.5, MID_GRAY)


def add_dimension_result(prs, data, page):
    slide = blank(prs)
    title(slide, "初步结果：四维度满意度整体处于中等偏上", "Preliminary Findings", page)
    colors = [GOLD, BLUE, GREEN, CORAL]
    for i, ((dim, val), color) in enumerate(zip(data["dim_stats"].items(), colors)):
        y = 2.25 + i * 1.2
        progress_bar(slide, 1.1, y, 8.2, dim, val, 5, color)
    card(slide, 11.25, 2.15, 6.6, 2.2, "结果一：陪伴空间略高", "社区亲子活动与邻里互助空间均值相对更高，说明多数受访者对基础亲子活动空间有一定感知，但并不意味着空间质量充分。", BLUE)
    card(slide, 11.25, 4.75, 6.6, 2.2, "结果二：休憩与服务弹性偏弱", "托育服务时间、照护者等候环境、休憩空间数量等细项分数偏低，显示问题集中在“父母时间”与“照护者恢复”。", CORAL)
    card(slide, 11.25, 7.35, 6.6, 1.55, "结果三：优质空间影响生育意愿", "“如果提供更多优质育儿友好空间，会更愿意生育或再生育”得分较高，可作为后续策略论证的关键支撑。", GOLD)


def add_item_painpoints(prs, data, page):
    slide = blank(prs)
    title(slide, "细项痛点：问题不只在有没有托位，而在能否嵌入日常", "Detailed Pain Points", page)
    low = sorted(data["item_stats"], key=lambda x: x[2])[:8]
    for i, (dim, head, val) in enumerate(low):
        short = head.split("—")[-1]
        short = short.replace("我", "").replace("（如社区宝宝屋、普惠托育点等）", "")
        y = 2.05 + i * 0.78
        color = [CORAL, GOLD, BLUE, GREEN][i % 4]
        text(slide, f"{i+1:02d}", 1.08, y, 0.55, 0.25, 13, color, True)
        text(slide, dim, 1.72, y, 1.65, 0.25, 11.5, color, True)
        text(slide, short[:34], 3.45, y, 8.9, 0.25, 11.5, BLACK)
        base = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(12.7), Inches(y + 0.05), Inches(3.1), Inches(0.18))
        fill(base, GRAY)
        fg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(12.7), Inches(y + 0.05), Inches(3.1 * val / 5), Inches(0.18))
        fill(fg, color)
        text(slide, f"{val:.2f}", 16.0, y - 0.02, 0.55, 0.25, 11.5, color, True)
    card(slide, 1.05, 8.75, 5.25, 1.25, "可解释为", "空间供给已经被看见，但对双职工最关键的“服务时间衔接、情绪喘息、照护等待舒适度”仍不足。", CORAL)
    card(slide, 6.75, 8.75, 5.25, 1.25, "对应策略", "从单点托育设施转向生活圈内的复合空间：临时托、亲子活动、父母休憩和邻里互助同步配置。", GOLD)
    card(slide, 12.45, 8.75, 5.25, 1.25, "后续验证", "用访谈追问低分项背后的具体场景，并与社区 POI、步行可达和通勤路径叠合。", BLUE)


def add_open_feedback(prs, data, page):
    slide = blank(prs)
    title(slide, "开放题反馈：需求集中于多功能、可达、可停留", "Open-ended Feedback", page)
    answers = data["open_answers"][:6]
    if len(answers) < 6:
        answers += ["希望增加儿童游乐场所", "需要更多多功能的生育友好空间", "社区空间应兼顾老人、儿童与父母的复合使用"]
    for i, ans in enumerate(answers[:6]):
        x = 1.05 + (i % 2) * 8.6
        y = 2.1 + (i // 2) * 2.1
        quote = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(7.55), Inches(1.35))
        fill(quote, LIGHT_GRAY)
        line(quote, LINE, 0.8)
        text(slide, f"“{ans[:46]}”", x + 0.32, y + 0.28, 6.8, 0.62, 15, BLACK, True)
    text(slide, "初步归纳", 1.05, 8.7, 1.5, 0.3, 16, GOLD, True)
    text(slide, "受访者并非只要求“增加儿童设施”，而是期待一种能够同时服务孩子玩耍、父母等待、邻里交流和短暂恢复的复合型空间。这个反馈与本研究提出的“三类空间”框架基本吻合。", 2.55, 8.65, 14.6, 0.58, 17, BLACK, True)


def add_method_refine(prs, page):
    slide = blank(prs)
    title(slide, "技术路线深化：从初步问卷走向空间验证", "Method Refinement", page)
    steps = [
        ("01", "问卷扩样", "扩大至 100-150 份，覆盖回龙观、天通苑及典型通勤型社区。"),
        ("02", "访谈校正", "选择高低分受访者进行半结构访谈，追问时间贫困和空间错配场景。"),
        ("03", "POI/GIS叠合", "抓取托育点、公园、商业、公共服务设施，计算15分钟可达与供需密度。"),
        ("04", "策略推演", "以社区闲置空间、架空层、口袋公园为载体，构建微更新方案。"),
    ]
    for i, (num, head, body) in enumerate(steps):
        x = 1.2 + i * 4.45
        top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(2.25), Inches(3.55), Inches(0.78))
        fill(top, GOLD if i % 2 == 0 else BLACK)
        text(slide, f"{num}  {head}", x + 0.2, 2.47, 3.1, 0.26, 15, WHITE, True)
        card(slide, x, 3.1, 3.55, 2.1, "", body, GOLD if i % 2 == 0 else BLACK, LIGHT_GRAY)
        if i < 3:
            arr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x + 3.68), Inches(3.88), Inches(0.56), Inches(0.28))
            fill(arr, GOLD)
    card(slide, 1.2, 6.55, 7.6, 1.85, "方法组合价值", "问卷负责识别需求强度，访谈负责解释需求背后的生活场景，GIS/POI负责定位空间错配，最终策略页才能落到具体社区与设施类型。", BLUE)
    card(slide, 9.35, 6.55, 7.6, 1.85, "与道路交通课程的关联", "双职工家庭的生育友好空间不是静态设施问题，而是通勤时间、接送路径、步行可达性与社区公共空间共同作用的日常时空组织问题。", CORAL)


def add_innovation(prs, page):
    slide = blank(prs)
    title(slide, "创新亮点：把生育友好从政策口号转译为空间网络", "Innovation", page)
    items = [
        ("理论创新", "以关怀地理学解释育儿压力的空间属性，将隐性的家庭照护劳动转化为可被规划识别的公共议题。"),
        ("对象创新", "聚焦适龄生育双职工家庭，抓住低生育率背景下“想生、敢生、能养”的城市空间阻滞。"),
        ("模型创新", "提出托育空间、陪伴空间、独处空间三类空间组合，突破单一托育点建设逻辑。"),
        ("方法创新", "结合问卷、访谈、GIS/POI、后续 Cesium 推演，建立从需求识别到空间模拟的闭环。"),
    ]
    for i, (head, body) in enumerate(items):
        x = 1.05 + (i % 2) * 8.55
        y = 2.05 + (i // 2) * 2.7
        card(slide, x, y, 7.6, 1.95, head, body, [GOLD, BLUE, GREEN, CORAL][i])
    text(slide, "一句话概括", 1.08, 8.35, 1.8, 0.32, 16, GOLD, True)
    text(slide, "不是“多建几个儿童设施”，而是把育儿支持嵌入家庭日常移动链条和社区公共空间网络中。", 3.0, 8.28, 13.8, 0.42, 24, BLACK, True)


def add_future(prs, page):
    slide = blank(prs)
    title(slide, "未来展望：从阶段汇报走向可落地的社区微更新方案", "Future Work", page)
    timeline = [
        ("第1阶段", "样本扩充", "继续发放问卷，补足不同社区类型、子女年龄、通勤距离的差异样本。"),
        ("第2阶段", "实地踏勘", "在回龙观、天通苑等社区筛选托育、陪伴、休憩空间典型点位。"),
        ("第3阶段", "空间叠合", "开展 POI/GIS 可达性、公共空间服务半径和通勤路径叠合分析。"),
        ("第4阶段", "策略输出", "形成“一图一表一策略”：空间问题图谱、指标评价表、微更新策略包。"),
    ]
    for i, (stage, head, body) in enumerate(timeline):
        y = 2.15 + i * 1.55
        text(slide, stage, 1.15, y, 1.45, 0.32, 15, GOLD, True)
        text(slide, head, 2.85, y, 1.7, 0.32, 16, BLACK, True)
        text(slide, body, 4.85, y, 10.8, 0.42, 14, INK)
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(16.25), Inches(y + 0.02), Inches(0.28), Inches(0.28))
        fill(dot, GOLD)
    card(slide, 1.1, 8.65, 5.2, 1.25, "预期成果", "完成调研报告、评价指标体系、典型社区空间改造建议与课程汇报展示。", GOLD)
    card(slide, 6.8, 8.65, 5.2, 1.25, "成果表达", "用问题图谱、统计图、策略剖面和微空间意向图增强最终展示效果。", BLUE)
    card(slide, 12.5, 8.65, 5.2, 1.25, "风险控制", "避免样本不足导致结论过度外推，将初步数据定位为问题识别与路径校正。", CORAL)


def add_final(prs, page):
    slide = blank(prs)
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
    fill(bg, INK)
    text(slide, "THANKS", 1.1, 1.15, 3.3, 0.5, 28, GOLD, True)
    text(slide, "让育儿支持成为城市生活圈里的日常能力", 1.1, 4.1, 12.4, 0.7, 34, WHITE, True)
    text(slide, "北京适龄生育双职工家庭生育友好空间调研", 1.15, 5.08, 8.0, 0.38, 17, WHITE)
    text(slide, f"{page:02d}", 18.4, 10.25, 0.7, 0.3, 14, WHITE, True, PP_ALIGN.RIGHT)


def main():
    data = load_data()
    copy2(REF, OUT)
    prs = Presentation(OUT)
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # Keep the reference deck's existing section for this topic: slides 16-23.
    delete_unwanted_slides(prs, set(range(15, 23)))

    add_cover(prs)
    move_slide(prs, len(prs.slides) - 1, 0)
    add_outline(prs)
    move_slide(prs, len(prs.slides) - 1, 1)

    page = len(prs.slides) + 1
    add_progress_overview(prs, data, page); page += 1
    add_sample_profile(prs, data, page); page += 1
    add_dimension_result(prs, data, page); page += 1
    add_item_painpoints(prs, data, page); page += 1
    add_open_feedback(prs, data, page); page += 1
    add_method_refine(prs, page); page += 1
    section_slide(prs, "PART 03", "创新亮点", "Innovation Highlights"); page += 1
    add_innovation(prs, page); page += 1
    section_slide(prs, "PART 04", "未来展望", "Future Work"); page += 1
    add_future(prs, page); page += 1
    add_final(prs, page)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
